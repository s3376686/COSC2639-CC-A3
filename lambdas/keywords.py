import json
import boto3
import botocore
from io import BytesIO
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.probability import FreqDist
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Ensure stopwords and punkt are downloaded (handle this in your deployment package)
nltk.data.path.append("/tmp")
nltk.download('punkt', download_dir='/tmp')
nltk.download('stopwords', download_dir='/tmp')

def lambda_handler(event, context):
    s3 = boto3.client('s3')

    # Extract bucket and key from the query parameters
    bucket = event.get('queryStringParameters', {}).get('bucket')
    key = event.get('queryStringParameters', {}).get('key')

    print(f"Extracted bucket: {bucket}")
    print(f"Extracted key: {key}")

    # Retrieve the file from S3
    try:
        response = s3.get_object(Bucket=bucket, Key=key)
        file_content = response['Body'].read()
    except botocore.exceptions.ClientError as e:
        error_code = e.response['Error']['Code']
        if error_code == 'NoSuchBucket':
            return {
                'statusCode': 400,
                'body': json.dumps(f'Bucket {bucket} does not exist')
            }
        elif error_code == 'NoSuchKey':
            return {
                'statusCode': 400,
                'body': json.dumps(f'Key {key} does not exist in bucket {bucket}')
            }
        else:
            raise e
    
    # Determine the file type and extract text accordingly
    if key.lower().endswith('.pdf'):
        all_text = extract_text_from_pdf(file_content)
    elif key.lower().endswith('.docx'):
        all_text = extract_text_from_docx(file_content)
    elif key.lower().endswith('.pptx'):
        all_text = extract_text_from_pptx(file_content)
    else:
        return {
            'statusCode': 400,
            'body': json.dumps('Unsupported file type.')
        }
    
    # Tokenize words and remove stopwords
    tokens = word_tokenize(all_text.lower())
    words = [word for word in tokens if word.isalpha()]  # Filter out punctuation and numbers
    stop_words = set(stopwords.words('english'))
    filtered_words = [word for word in words if word not in stop_words]
    
    # Get frequency distribution of words
    freq_dist = FreqDist(filtered_words)
    most_common_words = freq_dist.most_common(10)  # Adjust the number based on your needs

    return {
        'statusCode': 200,
        'body': json.dumps({
            'keywords': [word for word, freq in most_common_words]
        })
    }

def extract_text_from_pdf(file_content):
    pdf_reader = PdfReader(BytesIO(file_content))
    all_text = ""
    for page in pdf_reader.pages:
        all_text += page.extract_text()
    return all_text

def extract_text_from_docx(file_content):
    doc = Document(BytesIO(file_content))
    all_text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
    return all_text

def extract_text_from_pptx(file_content):
    prs = Presentation(BytesIO(file_content))
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                all_text += shape.text + "\n"
    return all_text